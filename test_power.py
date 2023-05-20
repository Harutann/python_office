from pptx import Presentation

prs = Presentation()
title_slide_layout = prs.slide_layouts[0]
slide = prs.slides.add_slide(title_slide_layout)
title = slide.shapes.title
subtitle = slide.placeholders[1]

title.text = "Hello, World!"
subtitle.text = "python-pptx was here!"

prs.save('test.pptx')

from docx import Document


def main():
    document = Document('sample.docx')
    document.add_picture('sample.jpg')
    count = 0
    for data in document.paragraphs:
        count += len(data.text)

    print(f'sample.docxの中の文字数は{count}個です。')

    document.save('sample_after.docx')


if __name__ == '__main__':
    main()